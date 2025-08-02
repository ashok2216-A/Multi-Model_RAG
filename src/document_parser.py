import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from gmft.auto import AutoTableDetector, AutoTableFormatter
from gmft.pdf_bindings import PyPDFium2Document

from src.config import DPI, MIN_PAGE_CHARS
from src.text_utils import uniquify_columns

# Initialize detectors
detector = AutoTableDetector()
formatter = AutoTableFormatter()

def is_page_searchable(page, min_chars=MIN_PAGE_CHARS):
    """Checks if a PDF page contains a minimum number of characters."""
    return len(page.get_text().strip()) >= min_chars

def process_ocr(page, dpi=DPI):
    """Performs OCR on a PDF page if it's not searchable."""
    pix = page.get_pixmap(dpi=dpi)
    ocr_pdf_bytes = pix.pdfocr_tobytes()
    return fitz.open("pdf", ocr_pdf_bytes)

def extract_tables_from_pdf_source(pdf_source, page_number):
    """Extracts tables from a given PDF source."""
    doc = PyPDFium2Document(pdf_source)
    tables = []
    try:
        page = doc[page_number]
        for cropped in detector.extract(page):
            formatted = formatter.extract(cropped, margin="auto", padding=None)
            df = formatted.df()
            tables.append((cropped.bbox, df))
    finally:
        doc.close()
    return tables

def extract_tables_from_fitz_doc(fitz_doc, page_number):
    """Extracts tables from a fitz document object."""
    single_page_doc = fitz.open()
    single_page_doc.insert_pdf(fitz_doc, from_page=page_number, to_page=page_number)
    pdf_bytes = single_page_doc.tobytes()
    single_page_doc.close()
    return extract_tables_from_pdf_source(pdf_bytes, 0)

def extract_text_and_tables_from_docx(path):
    """Extracts text and tables from a .docx file."""
    doc = Document(path)
    text = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    tables = []
    for table in doc.tables:
        table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append(table_data)
    return text, tables

def extract_text_and_tables_from_pptx(path):
    """Extracts text and tables from a .pptx file."""
    prs = Presentation(path)
    text, tables = [], []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
            if hasattr(shape, "has_table") and shape.has_table:
                table_data = [[cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells] for row in shape.table.rows]
                tables.append(table_data)
    return text, tables